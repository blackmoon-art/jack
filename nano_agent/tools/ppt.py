"""
PPT 生成工具 — 基于 python-pptx。

支持：标题页、内容页、项目符号页、双栏页。
自动安装 python-pptx（首次使用时）。
"""

import os
import subprocess
import sys


def _ensure_pptx():
    """确保 python-pptx 已安装。未安装时给出明确的安装指引。"""
    try:
        import pptx  # noqa: F401
        return
    except ImportError:
        raise ImportError(
            "python-pptx is not installed. Run: pip install python-pptx"
        )

class PPT:
    # 工具注册声明
    TOOLS = [
        ("create_ppt",
         "Generate a PowerPoint presentation (.pptx) with title slide and content slides. "
         "Supports title, content, bullets, and two-column layouts. Dark theme with purple accents. "
         "Auto-installs python-pptx on first use.",
         "create_ppt",
         {"title": {"type": "string", "description": "Main title of the presentation"},
          "slides": {"type": "array", "description": "List of slides. Each slide is an object with: type (title|content|bullets|two_column), title, body.",
                    "items": {"type": "object",
                              "properties": {"type": {"type": "string", "description": "Slide layout type: title, content, bullets, two_column"},
                                              "title": {"type": "string", "description": "Slide title"},
                                              "body": {"type": "string", "description": "Slide body text. For bullets type, use newlines to separate items."},
                                              "body_left": {"type": "string", "description": "Left column text (two_column type only)"},
                                              "body_right": {"type": "string", "description": "Right column text (two_column type only)"}}}},
          "filename": {"type": "string", "description": "Output filename (optional, defaults to title)"},
          "subtitle": {"type": "string", "description": "Subtitle for the title slide (optional)"}},
         ["title", "slides"]),
    ]

    def __init__(self, work_dir: str, charts_dir: str = ""):
        self.work_dir = work_dir
        self.charts_dir = charts_dir or work_dir

    def create_ppt(self, title: str, slides: list[dict], filename: str = "",
                   subtitle: str = "") -> str:
        """
        生成 PPT 文件。

        Args:
            title: PPT 主标题（也是第一页标题）
            slides: 幻灯片列表，每项是一个 dict:
                    - type: "title" | "content" | "bullets" | "two_column"
                    - title: 本页标题
                    - body: 正文内容（字符串，bullets 类型用换行分隔要点）
                    - body_left / body_right: 双栏内容
            filename: 输出文件名（不含路径，默认用标题）
            subtitle: 首页副标题

        Returns:
            生成结果信息
        """
        _ensure_pptx()

        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        # 16:9 宽屏
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # 颜色方案
        PRIMARY = RGBColor(0x7C, 0x3A, 0xED)   # 紫色
        DARK = RGBColor(0x1E, 0x1E, 0x2E)      # 深色
        LIGHT = RGBColor(0xE0, 0xE0, 0xE0)      # 浅色
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)

        def _add_text_box(slide, left, top, width, height, text,
                          font_size=18, color=LIGHT, bold=False, alignment=PP_ALIGN.LEFT):
            """添加文本框。"""
            txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                              Inches(width), Inches(height))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.bold = bold
            p.alignment = alignment
            return tf

        def _add_bullet_slide(slide, left, top, width, height, items,
                               font_size=16, color=LIGHT):
            """添加项目符号列表。"""
            txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                              Inches(width), Inches(height))
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, item in enumerate(items):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = item.strip()
                p.font.size = Pt(font_size)
                p.font.color.rgb = color
                p.space_after = Pt(8)
                # 项目符号
                p.level = 0
            return tf

        # ── 封面页 ──
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 背景色
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK

        # 主标题
        _add_text_box(slide, 1.5, 2.0, 10, 1.5, title,
                      font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        # 副标题
        if subtitle:
            _add_text_box(slide, 1.5, 3.8, 10, 1.0, subtitle,
                          font_size=20, color=PRIMARY, alignment=PP_ALIGN.CENTER)

        # ── 内容页 ──
        for i, slide_data in enumerate(slides):
            stype = slide_data.get("type", "content")
            stitle = slide_data.get("title", f"Slide {i+2}")
            body = slide_data.get("body", "")

            slide = prs.slides.add_slide(slide_layout)
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = DARK

            # 页面标题
            _add_text_box(slide, 0.8, 0.5, 11.5, 0.8, stitle,
                          font_size=32, color=PRIMARY, bold=True)

            # 标题下划线
            from pptx.util import Emu
            line = slide.shapes.add_shape(
                1,  # MSO_SHAPE.RECTANGLE
                Inches(0.8), Inches(1.4), Inches(11.5), Pt(3)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = PRIMARY
            line.line.fill.background()

            if stype == "title":
                # 大标题居中
                _add_text_box(slide, 1.5, 2.5, 10, 2.0, body,
                              font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

            elif stype == "bullets":
                # 项目符号
                items = [line for line in body.split("\n") if line.strip()]
                _add_bullet_slide(slide, 1.0, 2.0, 11, 5.0, items)

            elif stype == "two_column":
                body_left = slide_data.get("body_left", "")
                body_right = slide_data.get("body_right", "")
                _add_text_box(slide, 0.8, 2.0, 5.5, 5.0, body_left,
                              font_size=16, color=LIGHT)
                _add_text_box(slide, 6.8, 2.0, 5.5, 5.0, body_right,
                              font_size=16, color=LIGHT)

            else:  # content
                _add_text_box(slide, 0.8, 2.0, 11.5, 5.0, body,
                              font_size=18, color=LIGHT)

            # 页码
            _add_text_box(slide, 12.0, 7.0, 1.0, 0.4, str(i + 2),
                          font_size=10, color=RGBColor(0x66, 0x66, 0x66),
                          alignment=PP_ALIGN.RIGHT)

        # 保存
        if not filename:
            safe = "".join(c if c.isalnum() or c in "_-" else "_" for c in title)
            filename = f"{safe}.pptx"
        if not filename.endswith(".pptx"):
            filename += ".pptx"

        # 保存到 charts_dir（Web 可访问目录）
        if not filename:
            from datetime import datetime
            ts = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"ppt_{ts}.pptx"
        if not filename.endswith(".pptx"):
            filename += ".pptx"
        # 强制 ASCII 文件名，避免 URL 编码问题
        import re as _re
        import unicodedata as _ud
        safe = _ud.normalize('NFKD', filename).encode('ascii', 'ignore').decode('ascii')
        safe = _re.sub(r'[^a-zA-Z0-9._-]', '_', safe)
        if not safe.endswith('.pptx'):
            safe += '.pptx'
        if safe == '.pptx' or safe == '_.pptx':
            from datetime import datetime as _dt
            safe = f"ppt_{_dt.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        filename = safe
        filepath = os.path.join(self.charts_dir, filename)
        prs.save(filepath)
        url = f"/charts/{filename}"
        return f"PPT generated: {url} ({len(slides)+1} slides)\n[Download]({url})"

