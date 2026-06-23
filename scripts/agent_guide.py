#!/usr/bin/env python3
"""生成 Agent 开发教程 PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 颜色方案 ──
BLUE_DARK = RGBColor(0x1A, 0x27, 0x3A)
BLUE_ACCENT = RGBColor(0x29, 0x80, 0xB9)
TEAL = RGBColor(0x00, 0x96, 0x88)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)
CODE_BG = RGBColor(0x28, 0x2C, 0x34)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, font_size=16, color=DARK_TEXT, bullet=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.space_after = Pt(font_size * 0.5)
    return txBox

def add_icon_card(slide, left, top, width, height, icon_text, title, desc, bg_color=WHITE, title_color=BLUE_ACCENT):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    card.line.width = Pt(1)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), Inches(0.6), Inches(0.6),
                 icon_text, font_size=24, color=title_color, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.7), width - Inches(0.4), Inches(0.4),
                 title, font_size=14, color=title_color, bold=True)
    add_text_box(slide, left + Inches(0.2), top + Inches(1.1), width - Inches(0.4), height - Inches(1.3),
                 desc, font_size=11, color=DARK_TEXT)

# ============================================================
# 第1页：封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE_DARK)
add_shape(slide, Inches(0), Inches(0), Inches(0.3), Inches(7.5), TEAL)
add_shape(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06), TEAL)

add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.2),
             "🤖 AI Agent 开发实战教程", font_size=44, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(2.4), Inches(11), Inches(0.8),
             "从零到一构建智能体系统", font_size=24, color=TEAL)
add_text_box(slide, Inches(1), Inches(3.6), Inches(11), Inches(0.5),
             "涵盖：Agent 架构 · 工具调用 · 记忆系统 · 多智能体协作 · 生产部署", font_size=16, color=RGBColor(0xBB, 0xBB, 0xBB))
add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
             "2026 年 6 月", font_size=14, color=RGBColor(0x88, 0x88, 0x88))

# ============================================================
# 第2页：目录
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.3), Inches(7.5), BLUE_ACCENT)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             "📋 课程大纲", font_size=36, color=BLUE_DARK, bold=True)

toc_items = [
    ("1", "Agent 基本概念与原理"),
    ("2", "LLM 与 Agent 的核心能力"),
    ("3", "工具调用（Function Calling）"),
    ("4", "记忆系统构建"),
    ("5", "规划与推理策略"),
    ("6", "多智能体协作框架"),
    ("7", "Agent 安全与对齐"),
    ("8", "生产部署与监控"),
]

for i, (num, title) in enumerate(toc_items):
    y = Inches(1.6) + i * Inches(0.65)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y, Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BLUE_ACCENT if i % 2 == 0 else TEAL
    circle.line.fill.background()
    add_text_box(slide, Inches(1.2), y + Inches(0.05), Inches(0.45), Inches(0.4),
                 num, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.9), y + Inches(0.05), Inches(9), Inches(0.4),
                 title, font_size=18, color=DARK_TEXT)

# ============================================================
# 第3页：什么是 Agent
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE_ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "🧠 什么是 AI Agent？", font_size=32, color=BLUE_DARK, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.6),
             "AI Agent = 大语言模型（LLM）+ 感知 + 记忆 + 工具 + 规划 → 自主完成任务",
             font_size=18, color=BLUE_ACCENT, bold=True)

components = [
    ("🧠", "LLM 大脑", "理解意图、生成推理、\n决策规划"),
    ("👁️", "感知模块", "接收环境信息、\n解析用户输入、上下文理解"),
    ("💾", "记忆系统", "短期工作记忆、\n长期知识存储、经验回放"),
    ("🔧", "工具调用", "API/函数调用、\n数据库查询、代码执行"),
    ("🎯", "规划器", "任务分解、\n路径规划、反馈调整"),
]

for i, (icon, title, desc) in enumerate(components):
    x = Inches(0.6) + i * Inches(2.4)
    add_icon_card(slide, x, Inches(2.0), Inches(2.2), Inches(2.0), icon, title, desc)

add_shape(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(1.8), RGBColor(0xE8, 0xF4, 0xFD))
add_text_box(slide, Inches(1.2), Inches(4.7), Inches(10.5), Inches(0.4),
             "📌 Agent = LLM + 记忆 + 工具 + 规划能力", font_size=20, color=BLUE_ACCENT, bold=True)

lines = [
    "• 核心：LLM 作为\"大脑\"，负责推理与决策",
    "• 增强：通过工具调用突破 LLM 知识截止日期和能力边界",
    "• 闭环：观察 → 思考 → 行动 → 反馈 → 调整（ReAct 范式）",
]
add_multi_text(slide, Inches(1.2), Inches(5.2), Inches(10.5), Inches(1.0), lines, font_size=14, color=DARK_TEXT)

# ============================================================
# 第4页：Agent 发展历史
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE_ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "📅 Agent 发展历程", font_size=32, color=BLUE_DARK, bold=True)

milestones = [
    ("2022", "ChatGPT 发布\nLLM 能力引爆 Agent 研究"),
    ("2023", "AutoGPT / AgentGPT\n掀起自主 Agent 热潮"),
    ("2023", "ReAct / CoT 推理框架\nChain-of-Thought 提示工程"),
    ("2024", "LangChain / CrewAI / AutoGen\n框架成熟，生产可用"),
    ("2024", "MCP 协议标准\n工具调用标准化"),
    ("2025", "多模态 Agent\n视觉+语音+代码全栈"),
    ("2026", "Agent 工程化\n企业级部署、安全合规"),
]

for i, (year, desc) in enumerate(milestones):
    x = Inches(0.5) + i * Inches(1.8)
    if i < len(milestones) - 1:
        line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.25) + i * Inches(1.8), Inches(2.3), Inches(0.04), Inches(1.5))
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        line_shape.line.fill.background()
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.1) + i * Inches(1.8), Inches(2.1), Inches(0.3), Inches(0.3))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BLUE_ACCENT if i < 3 else (ORANGE if i == 3 else TEAL)
    circle.line.fill.background()
    add_text_box(slide, x, Inches(1.4), Inches(1.6), Inches(0.6),
                 year, font_size=16, color=BLUE_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(3.8), Inches(1.6), Inches(1.8),
                 desc, font_size=11, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

# ============================================================
# 第5页：LLM 与 Agent 核心能力
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE_ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "⚡ LLM 驱动的 Agent 核心能力", font_size=32, color=BLUE_DARK, bold=True)

capabilities = [
    ("🔍", "推理能力", "ReAct / CoT / ToT\n逻辑推理、数学计算\n因果关系分析"),
    ("📝", "内容生成", "写作、代码生成\n报告总结、翻译\n创意内容创作"),
    ("🔗", "工具调用", "Function Calling\nAPI 集成\n数据库/代码执行"),
    ("🧩", "多步规划", "任务分解（HuggingGPT）\n子任务调度\n动态重规划"),
    ("💬", "多轮对话", "上下文维护\n状态跟踪\n长期记忆检索"),
    ("🛡️", "安全对齐", "指令遵循\n安全过滤\n红队测试防护"),
]

for i, (icon, title, desc) in enumerate(capabilities):
    row = i // 3
    col = i % 3
    x = Inches(0.6) + col * Inches(4.1)
    y = Inches(1.3) + row * Inches(2.8)
    add_icon_card(slide, x, y, Inches(3.8), Inches(2.4), icon, title, desc)

# ============================================================
# 第6页：工具调用 Function Calling
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE_ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "🔧 工具调用（Function Calling）", font_size=32, color=BLUE_DARK, bold=True)

add_shape(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(4.5), RGBColor(0xE8, 0xF4, 0xFD))
add_text_box(slide, Inches(1.0), Inches(1.4), Inches(5), Inches(0.4),
             "工作流程", font_size=18, color=BLUE_ACCENT, bold=True)

steps = [
    '1️⃣ 用户提问：帮我查一下北京的天气',
    '2️⃣ LLM 分析意图 → 选择工具：get_weather(city=\"北京\")',
    '3️⃣ 解析函数参数，调用外部 API',
    '4️⃣ 返回结果：北京今日 25°C，晴',
    '5️⃣ LLM 组织自然语言回复给用户',
]
add_multi_text(slide, Inches(1.2), Inches(1.9), Inches(5), Inches(3.5), steps, font_size=14, color=DARK_TEXT)

code_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.3), Inches(5.8), Inches(4.5))
code_shape.fill.solid()
code_shape.fill.fore_color.rgb = CODE_BG
code_shape.line.fill.background()

code_text = """tools = [
  {
    "type": "function",
    "function": {
      "name": "get_weather",
      "description": "获取天气",
      "parameters": {
        "type": "object",
        "properties": {
          "city": {
            "type": "string",
            "description": "城市名"
          }
        },
        "required": ["city"]
      }
    }
  }
]

# LLM 返回工具调用
response = client.chat.completions.create(
    model="gpt-4",
    messages=messages,
    tools=tools
)
# 执行工具并返回结果"""

add_text_box(slide, Inches(7.0), Inches(1.5), Inches(5.4), Inches(4.0),
             code_text, font_size=11, color=RGBColor(0xE0, 0xE0, 0xE0), font_name="Courier New")

add_shape(slide, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.6), RGBColor(0xFF, 0xF3, 0xE0))
add_text_box(slide, Inches(1.0), Inches(6.2), Inches(11), Inches(0.4),
             "💡 主流框架: OpenAI Function Calling · Anthropic Tool Use · Google Vertex AI · LangChain Tools · MCP Protocol",
             font_size=13, color=ORANGE)

# ============================================================
# 第7页：记忆系统
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE_ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "💾 记忆系统架构", font_size=32, color=BLUE_DARK, bold=True)

memories = [
    ("🔄", "短期记忆", "上下文窗口（几K~200K tokens）\n多轮对话历史", BLUE_ACCENT),
    ("📦", "长期记忆", "向量数据库（Chroma/Pinecone）\nRAG 检索增强", TEAL),
    ("⚡", "工作记忆", "任务状态追踪\n中间推理结果、待办事项", ORANGE),
]

for i, (icon, title, desc, color) in enumerate(memories):
    x = Inches(0.6) + i * Inches(4.1)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.3), Inches(3.8), Inches(2.5))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
    card.line.color.rgb = color
    card.line.width = Pt(2)
    add_text_box(slide, x + Inches(0.2), Inches(1.4), Inches(3.4), Inches(0.5),
                 f"{icon}  {title}", font_size=18, color=color, bold=True)
    add_text_box(slide, x + Inches(0.2), Inches(2.0), Inches(3.4), Inches(1.5),
                 desc, font_size=13, color=DARK_TEXT)

add_shape(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(2.8), RGBColor(0xF5, 0xF5, 0xF5))
add_text_box(slide, Inches(1.0), Inches(4.3), Inches(5), Inches(0.4),
             "📎 RAG 检索增强生成", font_size=18, color=BLUE_ACCENT, bold=True)

rag_steps = [
    "① 用户提问 → 嵌入向量化",
    "② 向量数据库语义检索 Top-K",
    "③ 检索结果 + 原始问题 → LLM",
    "④ LLM 基于上下文生成回答",
]
add_multi_text(slide, Inches(1.0), Inches(4.8), Inches(5.5), Inches(2.0), rag_steps, font_size=14, color=DARK_TEXT)

tech_text = """技术选型推荐

向量数据库：Chroma / Pinecone / Weaviate
嵌入模型：text-embedding-3-small / BGE
长期存储：PostgreSQL + pgvector
缓存方案：Redis + 语义缓存
混合检索：BM25 + 向量检索"""
add_text_box(slide, Inches(6.5), Inches(4.8), Inches(5.5), Inches(2.2),
             tech_text, font_size=13, color=DARK_TEXT)

# ============================================================
# 第8页：规划与推理策略
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE_ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "🎯 规划与推理策略", font_size=32, color=BLUE_DARK, bold=True)

strategies = [
    ("ReAct", "推理+行动交替\nThought → Action → Observation\n循环迭代直至任务完成", "2023"),
    ("CoT", "链式思考提示\n将复杂问题分解为\n中间推理步骤", "2022"),
    ("ToT", "思维树\n多路径搜索+剪枝\n探索多种推理路径", "2023"),
    ("Plan&Exec", "先规划再执行\n分解为子任务\n依次调度工具", "2024"),
]

for i, (name, desc, year) in enumerate(strategies):
    x = Inches(0.5) + i * Inches(3.1)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.3), Inches(2.9), Inches(2.8))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BLUE_ACCENT
    card.line.width = Pt(1.5)
    add_text_box(slide, x + Inches(0.2), Inches(1.4), Inches(2.5), Inches(0.4),
                 name, font_size=22, color=BLUE_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), Inches(2.0), Inches(2.6), Inches(1.5),
                 desc, font_size=12, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), Inches(3.4), Inches(2.6), Inches(0.3),
                 f"({year})", font_size=10, color=RGBColor(0x99, 0x99, 0x99), alignment=PP_ALIGN.CENTER)

code_shape2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.5))
code_shape2.fill.solid()
code_shape2.fill.fore_color.rgb = CODE_BG
code_shape2.line.fill.background()

react_code = """# ReAct 循环伪代码
def agent_loop(task):
    while not done:
        # Thought: 推理当前状态
        thought = llm.reason(f"当前任务: {task}\\n已有信息: {memory}")
        # Action: 选择并执行工具
        action = llm.choose_action(thought, tools)
        observation = execute_tool(action)
        # 更新记忆
        memory.add(f"Thought: {thought}\\nAction: {action}\\nObs: {observation}")
        if llm.is_complete(task, memory):
            return llm.summarize(memory)"""

add_text_box(slide, Inches(1.0), Inches(4.6), Inches(11), Inches(2.2),
             react_code, font_size=11, color=RGBColor(0xE0, 0xE0, 0xE0), font_name="Courier New")

# ============================================================
# 第9页：多智能体协作
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE_ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "👥 多智能体协作框架", font_size=32, color=BLUE_DARK, bold=True)

patterns = [
    ("🏛️", "层级式", "Manager Agent 分配任务\nWorker Agent 执行\n结果汇总上报", BLUE_ACCENT),
    ("🔄", "流水线式", "每个 Agent 负责一个环节\n输出作为下个 Agent 输入\n类似工厂生产线", TEAL),
    ("🤝", "协商式", "Agent 间自由讨论\n投票/辩论达成共识\n适用于复杂决策", ORANGE),
    ("🌐", "市场式", "Agent 作为独立服务\n通过消息队列通信\n事件驱动松耦合", RGBColor(0x8E, 0x44, 0xAD)),
]

for i, (icon, title, desc, color) in enumerate(patterns):
    x = Inches(0.4) + i * Inches(3.15)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.3), Inches(2.95), Inches(2.3))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
    card.line.color.rgb = color
    card.line.width = Pt(2)
    add_text_box(slide, x + Inches(0.15), Inches(1.4), Inches(2.6), Inches(0.4),
                 f"{icon} {title}", font_size=16, color=color, bold=True)
    add_text_box(slide, x + Inches(0.15), Inches(1.9), Inches(2.6), Inches(1.5),
                 desc, font_size=12, color=DARK_TEXT)

add_shape(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(3.0), RGBColor(0xF5, 0xF5, 0xF5))
add_text_box(slide, Inches(1.0), Inches(4.1), Inches(5), Inches(0.4),
             "📊 主流多 Agent 框架对比", font_size=18, color=BLUE_ACCENT, bold=True)

frameworks = [
    ("LangGraph", "LangChain", "有向图编排, 状态持久化", "⭐⭐⭐⭐⭐"),
    ("CrewAI", "独立", "角色分工, 任务委派", "⭐⭐⭐⭐"),
    ("AutoGen", "Microsoft", "对话驱动, 灵活配置", "⭐⭐⭐⭐"),
    ("OpenAI Swarm", "OpenAI", "轻量级, 函数调用", "⭐⭐⭐"),
]

headers = ["框架", "开发者", "核心特性", "成熟度"]
col_widths = [Inches(2.5), Inches(2.0), Inches(4.5), Inches(2.0)]
for j, h in enumerate(headers):
    x_pos = Inches(1.0)
    for k in range(j):
        x_pos += col_widths[k]
    add_text_box(slide, x_pos, Inches(4.6), col_widths[j], Inches(0.4),
                 h, font_size=12, color=BLUE_ACCENT, bold=True)

for i, fw in enumerate(frameworks):
    y = Inches(5.1) + i * Inches(0.45)
    x_pos = Inches(1.0)
    for j, val in enumerate(fw):
        add_text_box(slide, x_pos, y, col_widths[j], Inches(0.4),
                     val, font_size=11, color=DARK_TEXT)
        if j < len(fw) - 1:
            x_pos += col_widths[j]

# ============================================================
# 第10页：Agent 安全与对齐
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE_ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "🛡️ Agent 安全与对齐", font_size=32, color=BLUE_DARK, bold=True)

risks = [
    ("⚠️ 提示注入", "恶意输入劫持 Agent 行为\n解决方案：输入过滤 + 权限隔离"),
    ("🔓 工具滥用", "Agent 调用敏感 API\n解决方案：白名单 + 人工审批"),
    ("📋 幻觉控制", "生成虚假信息\n解决方案：RAG 事实校验 + 溯源"),
    ("🔒 数据泄露", "记忆泄露敏感信息\n解决方案：数据脱敏 + 访问控制"),
    ("🎯 对齐失败", "Agent 行为偏离预期\n解决方案：RLHF + 约束提示"),
]

for i, (title, desc) in enumerate(risks):
    y = Inches(1.3) + i * Inches(1.05)
    bg_c = RGBColor(0xFD, 0xED, 0xED) if i < 2 else RGBColor(0xFE, 0xF5, 0xE7)
    add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(0.9), bg_c)
    add_text_box(slide, Inches(1.0), y + Inches(0.05), Inches(3), Inches(0.4),
                 title, font_size=15, color=BLUE_ACCENT, bold=True)
    add_text_box(slide, Inches(4.0), y + Inches(0.05), Inches(8), Inches(0.7),
                 desc, font_size=13, color=DARK_TEXT)

add_shape(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.8), RGBColor(0xE8, 0xF4, 0xFD))
add_text_box(slide, Inches(1.0), Inches(6.1), Inches(11), Inches(0.5),
             "✅ 最佳实践：最小权限原则 · 工具调用审计日志 · 人工在环（HITL） · 定期红队测试 · 输出内容过滤",
             font_size=14, color=BLUE_ACCENT, bold=True)

# ============================================================
# 第11页：生产部署
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE_ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "🚀 生产部署与监控", font_size=32, color=BLUE_DARK, bold=True)

deploy_icons = ["🌐", "⚡", "📊", "🔁", "✅"]
deploy_items = [
    ("API 服务", "FastAPI / Flask\nOpenAI 兼容接口\n异步流式响应"),
    ("缓存", "Redis 语义缓存\n减少 API 调用\n降低延迟"),
    ("可观测性", "LangSmith / Arize\n链路追踪\nToken 用量监控"),
    ("弹性", "自动扩缩容\n限流降级\n故障转移"),
    ("评估", "离线评估集\n在线 A/B 测试\n用户满意度"),
]

for i, (title, desc) in enumerate(deploy_items):
    x = Inches(0.4) + i * Inches(2.5)
    add_icon_card(slide, x, Inches(1.3), Inches(2.3), Inches(2.0),
                  deploy_icons[i], title, desc)

add_shape(slide, Inches(0.8), Inches(3.8), Inches(11.5), Inches(3.0), RGBColor(0xF5, 0xF5, 0xF5))
add_text_box(slide, Inches(1.0), Inches(3.9), Inches(5), Inches(0.4),
             "🏗️ 生产架构推荐", font_size=18, color=BLUE_ACCENT, bold=True)

arch_text = """用户 → Nginx/ALB → API Gateway → Agent Service (K8s) 
                                        ├── LLM (OpenAI/Claude/本地)
                                        ├── 工具执行器 (Sandbox)
                                        ├── 向量数据库 (Chroma/Pinecone)
                                        └── 消息队列 (Redis/Kafka)

监控栈: Prometheus + Grafana | 日志: ELK | 链路: LangSmith | 告警: PagerDuty"""
add_text_box(slide, Inches(1.0), Inches(4.4), Inches(11), Inches(2.0),
             arch_text, font_size=12, color=DARK_TEXT)

# ============================================================
# 第12页：动手实践项目
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE_ACCENT)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
             "💻 动手实践：构建你的第一个 Agent", font_size=32, color=BLUE_DARK, bold=True)

projects = [
    ("📚 项目1：智能客服 Agent", "• 对接 OpenAI API\n• 配置知识库 RAG\n• 实现上下文记忆\n• 支持多轮对话", BLUE_ACCENT),
    ("📊 项目2：数据分析 Agent", "• 连接数据库\n• 自然语言转 SQL\n• 自动生成图表\n• 报告自动总结", TEAL),
    ("🤖 项目3：多 Agent 协作", "• 使用 CrewAI\n• 定义角色和任务\n• Agent 间通信\n• 结果汇总输出", ORANGE),
]

for i, (title, desc, color) in enumerate(projects):
    x = Inches(0.5) + i * Inches(4.15)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.3), Inches(3.9), Inches(2.6))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
    card.line.color.rgb = color
    card.line.width = Pt(2)
    add_text_box(slide, x + Inches(0.2), Inches(1.4), Inches(3.5), Inches(0.4),
                 title, font_size=16, color=color, bold=True)
    add_text_box(slide, x + Inches(0.2), Inches(1.9), Inches(3.5), Inches(1.8),
                 desc, font_size=13, color=DARK_TEXT)

add_shape(slide, Inches(0.8), Inches(4.3), Inches(11.5), Inches(2.7), RGBColor(0xE8, 0xF4, 0xFD))
add_text_box(slide, Inches(1.0), Inches(4.4), Inches(5), Inches(0.4),
             "📖 推荐学习资源", font_size=18, color=BLUE_ACCENT, bold=True)

resources = [
    "📘 官方文档: OpenAI Agents · Anthropic Tool Use · LangChain",
    "📗 开源项目: AutoGPT, MetaGPT, ChatDev, OpenDevin",
    "📙 论文: ReAct (Yao et al.), Toolformer (Schick et al.), Tree of Thoughts (Yao et al.)",
    "🎓 在线课程: DeepLearning.AI Building Agentic RAG, LangChain 官方教程",
    "💬 社区: Discord LangChain, Twitter/X #AgenticAI",
]
add_multi_text(slide, Inches(1.0), Inches(4.9), Inches(11), Inches(1.8), resources, font_size=13, color=DARK_TEXT)

# ============================================================
# 第13页：结语
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE_DARK)
add_shape(slide, Inches(0), Inches(0), Inches(0.3), Inches(7.5), TEAL)
add_shape(slide, Inches(0), Inches(3.5), Inches(13.333), Inches(0.04), TEAL)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.0),
             "🎯 总结与展望", font_size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.5),
             "Agent 是 AI 从「对话」走向「行动」的关键一步", font_size=20, color=TEAL)

summary_lines = [
    "• Agent = LLM + 记忆 + 工具 + 规划",
    "• 单 Agent 解决独立任务，多 Agent 协作处理复杂流程",
    "• 安全、可观测性、可靠性是生产落地的关键",
    "• 2026 年 Agent 正从实验走向工程化",
]
add_multi_text(slide, Inches(1), Inches(3.8), Inches(11), Inches(2.0), summary_lines, font_size=16, color=RGBColor(0xDD, 0xDD, 0xDD))

add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.6),
             "🚀 未来趋势：多模态 Agent · Agent 市场 · 自主 Agent 团队 · Agent 即服务 (AaaS)",
             font_size=15, color=RGBColor(0x88, 0xDD, 0xCC))

add_text_box(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.5),
             "Thanks! 欢迎交流 🙌", font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# ── 保存 ──
output_path = "agent_guide.ppt"
prs.save(output_path)
print(f"✅ PPT has been saved to: {output_path}")
