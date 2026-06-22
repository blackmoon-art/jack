from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 颜色方案
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_CARD = RGBColor(0x16, 0x21, 0x3E)
ACCENT_BLUE = RGBColor(0x00, 0xB4, 0xD8)
ACCENT_CYAN = RGBColor(0x00, 0xE5, 0xFF)
ACCENT_PURPLE = RGBColor(0x7B, 0x2C, 0xBF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x76)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
YELLOW = RGBColor(0xFF, 0xD7, 0x00)

def add_bg(slide, color=BG_DARK):
    """设置纯色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, alpha=None):
    """添加矩形色块"""
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Arial"):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_text(slide, left, top, width, height, items, font_size=18, color=LIGHT_GRAY, spacing=Pt(8)):
    """添加列表文本"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"▸ {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.space_after = spacing
    return txBox

def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    """添加装饰线"""
    shape = add_shape(slide, left, top, width, Pt(4), color)
    return shape

# ===== 第1页：封面 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
add_bg(slide)

# 顶部装饰条
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.15), ACCENT_BLUE)

# 标题
add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
             "AI Agent 开发", font_size=54, color=WHITE, bold=True)

# 副标题
add_text_box(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(0.8),
             "从入门到实战：构建智能代理系统", font_size=28, color=ACCENT_CYAN)

# 装饰线
add_accent_line(slide, Inches(1.5), Inches(3.9), Inches(3), ACCENT_BLUE)

# 底部信息
add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
             "2026 · Agent 开发实践指南", font_size=16, color=LIGHT_GRAY)

# 渐变装饰块 - 右侧
add_shape(slide, Inches(11), Inches(1.5), Inches(0.08), Inches(4.5), ACCENT_PURPLE)
add_shape(slide, Inches(11.2), Inches(2.0), Inches(0.06), Inches(4.0), ACCENT_BLUE)

# ===== 第2页：目录 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(1.5), ACCENT_CYAN)
add_text_box(slide, Inches(0.8), Inches(1.1), Inches(5), Inches(0.6),
             "目 录", font_size=36, color=WHITE, bold=True)

items = [
    ("01", "什么是 AI Agent", ACCENT_CYAN),
    ("02", "Agent 核心架构", ACCENT_BLUE),
    ("03", "工具调用与 Function Calling", ACCENT_GREEN),
    ("04", "记忆与上下文管理", ACCENT_PURPLE),
    ("05", "多 Agent 协作", YELLOW),
    ("06", "实战项目：智能助手", ACCENT_CYAN),
    ("07", "最佳实践与挑战", ACCENT_BLUE),
]

for i, (num, title, color) in enumerate(items):
    y = Inches(2.2) + Inches(i * 0.7)
    add_text_box(slide, Inches(1.2), y, Inches(0.6), Inches(0.5),
                 num, font_size=24, color=color, bold=True)
    add_text_box(slide, Inches(2.0), y, Inches(8), Inches(0.5),
                 title, font_size=22, color=WHITE)

# ===== 第3页：什么是 AI Agent =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.2), ACCENT_CYAN)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(5), Inches(0.6),
             "什么是 AI Agent", font_size=32, color=WHITE, bold=True)

# 左侧定义
add_shape(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(2.8), ACCENT_CARD := RGBColor(0x1E, 0x2A, 0x4A))
add_text_box(slide, Inches(1.2), Inches(1.9), Inches(5), Inches(0.5),
             "🤖 Agent 定义", font_size=22, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(1.2), Inches(2.5), Inches(5), Inches(1.5),
             "AI Agent 是一个能够自主感知环境、制定计划、\n使用工具并执行行动的智能系统。它不仅仅是\n一个对话模型，而是一个具备自主决策能力的\n智能体。",
             font_size=16, color=LIGHT_GRAY)

# 右侧关键特征
add_shape(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(2.8), ACCENT_CARD)
add_text_box(slide, Inches(7.4), Inches(1.9), Inches(4.8), Inches(0.5),
             "⚡ 关键特征", font_size=22, color=ACCENT_GREEN, bold=True)
add_bullet_text(slide, Inches(7.4), Inches(2.5), Inches(4.8), Inches(1.8),
                ["自主性：无需人类持续干预", "感知力：理解环境与任务上下文",
                 "行动力：调用工具、执行操作", "学习力：从反馈中持续优化"],
                font_size=16, color=LIGHT_GRAY, spacing=Pt(6))

# 底部对比
add_shape(slide, Inches(0.8), Inches(4.9), Inches(11.7), Inches(2.0), ACCENT_CARD)
add_text_box(slide, Inches(1.2), Inches(5.1), Inches(11), Inches(0.4),
             "📊 传统 LLM vs AI Agent", font_size=20, color=YELLOW, bold=True)
add_bullet_text(slide, Inches(1.2), Inches(5.6), Inches(5), Inches(1.0),
                ["传统 LLM：一问一答，被动响应"],
                font_size=16, color=LIGHT_GRAY, spacing=Pt(4))
add_bullet_text(slide, Inches(6.5), Inches(5.6), Inches(5.5), Inches(1.0),
                ["AI Agent：目标驱动 → 规划 → 执行 → 观察 → 迭代"],
                font_size=16, color=ACCENT_CYAN, spacing=Pt(4))

# ===== 第4页：Agent 核心架构 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.2), ACCENT_CYAN)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(8), Inches(0.6),
             "Agent 核心架构", font_size=32, color=WHITE, bold=True)

# 架构图 - 5个模块
modules = [
    ("🧠 感知模块", "接收并理解\n用户输入与\n环境状态", ACCENT_CYAN),
    ("💭 推理规划", "分析任务\n制定执行\n计划步骤", ACCENT_BLUE),
    ("🔧 工具调用", "调用API/\n代码执行/\n数据库操作", ACCENT_GREEN),
    ("📝 记忆系统", "短期/长期\n记忆管理\n上下文存储", ACCENT_PURPLE),
    ("⚡ 执行引擎", "协调各模块\n驱动整个\n循环流程", YELLOW),
]

for i, (title, desc, color) in enumerate(modules):
    x = Inches(0.6) + Inches(i * 2.5)
    # 模块卡片
    add_shape(slide, x, Inches(1.8), Inches(2.2), Inches(2.8), ACCENT_CARD)
    # 顶部色条
    add_shape(slide, x, Inches(1.8), Inches(2.2), Inches(0.08), color)
    # 模块名
    add_text_box(slide, x + Inches(0.2), Inches(2.0), Inches(1.8), Inches(0.5),
                 title, font_size=16, color=color, bold=True)
    # 描述
    add_text_box(slide, x + Inches(0.2), Inches(2.7), Inches(1.8), Inches(1.5),
                 desc, font_size=14, color=LIGHT_GRAY)
    # 箭头（最后一个不画）
    if i < len(modules) - 1:
        add_text_box(slide, x + Inches(2.1), Inches(2.8), Inches(0.5), Inches(0.5),
                     "→", font_size=28, color=ACCENT_CYAN, bold=True)

# 底部说明
add_shape(slide, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.8), ACCENT_CARD)
add_text_box(slide, Inches(1.0), Inches(5.2), Inches(11.3), Inches(0.4),
             "🔄 Agent 工作循环", font_size=20, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(1.0), Inches(5.7), Inches(11.3), Inches(0.8),
             "感知 → 规划 → 执行(工具调用) → 观察结果 → 更新记忆 → 再次规划 → 直至完成目标",
             font_size=16, color=YELLOW)

# ===== 第5页：工具调用与 Function Calling =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.2), ACCENT_CYAN)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(8), Inches(0.6),
             "工具调用与 Function Calling", font_size=32, color=WHITE, bold=True)

# 左侧：概念
add_shape(slide, Inches(0.8), Inches(1.6), Inches(5.8), Inches(5.0), ACCENT_CARD)
add_text_box(slide, Inches(1.2), Inches(1.8), Inches(5), Inches(0.4),
             "🔌 Function Calling 流程", font_size=22, color=ACCENT_BLUE, bold=True)

add_text_box(slide, Inches(1.2), Inches(2.4), Inches(5), Inches(3.8),
             "1. 定义工具函数（描述 + 参数 schema）\n\n"
             "2. LLM 分析用户意图 → 决定调用哪些工具\n\n"
             "3. 返回结构化 JSON（函数名 + 参数）\n\n"
             "4. 本地执行函数，获取结果\n\n"
             "5. 将结果送回 LLM，生成最终回复",
             font_size=15, color=LIGHT_GRAY)

# 右侧：代码示例
add_shape(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(5.0), ACCENT_CARD)
add_text_box(slide, Inches(7.4), Inches(1.8), Inches(4.8), Inches(0.4),
             "📝 示例：天气查询 Agent", font_size=22, color=ACCENT_GREEN, bold=True)

# 用文本模拟代码块
code_text = (
    'tools = [{\n'
    '  "type": "function",\n'
    '  "function": {\n'
    '    "name": "get_weather",\n'
    '    "description": "查询天气",\n'
    '    "parameters": {\n'
    '      "type": "object",\n'
    '      "properties": {\n'
    '        "city": {"type": "string"},\n'
    '        "date": {"type": "string"}\n'
    '      },\n'
    '      "required": ["city"]\n'
    '    }\n'
    '  }\n'
    '}]\n\n'
    "# Agent 自动调用\n"
    "response = client.chat(\n"
    '  model="...",\n'
    "  tools=tools,\n"
    "  messages=[...]\n"
    ")"
)
add_text_box(slide, Inches(7.4), Inches(2.4), Inches(4.8), Inches(3.8),
             code_text, font_size=13, color=ACCENT_CYAN)

# ===== 第6页：记忆系统 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.2), ACCENT_CYAN)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(8), Inches(0.6),
             "记忆与上下文管理", font_size=32, color=WHITE, bold=True)

# 三栏记忆类型
memories = [
    ("🧠 短期记忆", "• 当前对话上下文\n• Token 窗口管理\n• 滑动窗口策略\n• 摘要压缩技术", ACCENT_CYAN),
    ("📚 长期记忆", "• 向量数据库存储\n• 语义检索（RAG）\n• 知识图谱关联\n• 用户偏好持久化", ACCENT_GREEN),
    ("⚡ 工作记忆", "• 当前任务状态\n• 中间推理步骤\n• 待办事项跟踪\n• 执行进度管理", ACCENT_PURPLE),
]

for i, (title, content, color) in enumerate(memories):
    x = Inches(0.8) + Inches(i * 4.1)
    add_shape(slide, x, Inches(1.6), Inches(3.7), Inches(4.0), ACCENT_CARD)
    add_shape(slide, x, Inches(1.6), Inches(3.7), Inches(0.08), color)
    add_text_box(slide, x + Inches(0.3), Inches(1.9), Inches(3.1), Inches(0.5),
                 title, font_size=20, color=color, bold=True)
    add_text_box(slide, x + Inches(0.3), Inches(2.6), Inches(3.1), Inches(2.5),
                 content, font_size=14, color=LIGHT_GRAY)

# 底部架构
add_shape(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(1.0), ACCENT_CARD)
add_text_box(slide, Inches(1.2), Inches(6.1), Inches(11), Inches(0.6),
             "💡 最佳实践：短期记忆处理实时对话，长期记忆通过 RAG 检索增强，工作记忆跟踪多步骤任务状态",
             font_size=15, color=YELLOW)

# ===== 第7页：多 Agent 协作 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.2), ACCENT_CYAN)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(8), Inches(0.6),
             "多 Agent 协作", font_size=32, color=WHITE, bold=True)

# 架构图卡片
patterns = [
    ("👑 主管- Worker 模式", "一个主管 Agent 分配\n任务给多个 Worker\nAgent，汇总结果", ACCENT_CYAN),
    ("🤝 对话式协作", "多个 Agent 通过\n结构化对话协商\n解决问题", ACCENT_GREEN),
    ("🏗️ 流水线模式", "每个 Agent 负责\n一个阶段，输出\n作为下一阶段输入", ACCENT_BLUE),
    ("🌐 市场模式", "Agent 作为服务提供者\n通过竞标/协商\n完成任务", ACCENT_PURPLE),
]

for i, (title, desc, color) in enumerate(patterns):
    x = Inches(0.5) + Inches(i * 3.15)
    add_shape(slide, x, Inches(1.6), Inches(2.9), Inches(2.5), ACCENT_CARD)
    add_shape(slide, x, Inches(1.6), Inches(2.9), Inches(0.08), color)
    add_text_box(slide, x + Inches(0.2), Inches(1.9), Inches(2.5), Inches(0.5),
                 title, font_size=15, color=color, bold=True)
    add_text_box(slide, x + Inches(0.2), Inches(2.5), Inches(2.5), Inches(1.2),
                 desc, font_size=14, color=LIGHT_GRAY)

# 使用场景
add_shape(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5), ACCENT_CARD)
add_text_box(slide, Inches(0.9), Inches(4.7), Inches(5), Inches(0.4),
             "🎯 典型应用场景", font_size=20, color=YELLOW, bold=True)

scenarios = [
    ("软件工程", "产品经理 Agent +\n开发 Agent + QA Agent"),
    ("客服系统", "分流 Agent + 业务 Agent\n+ 升级 Agent"),
    ("数据分析", "数据获取 Agent +\n分析 Agent + 可视化 Agent"),
    ("自动化运维", "监控 Agent + 诊断 Agent\n+ 修复 Agent"),
]

for i, (title, desc) in enumerate(scenarios):
    x = Inches(0.9) + Inches(i * 3.05)
    add_text_box(slide, x, Inches(5.2), Inches(2.8), Inches(0.4),
                 f"▸ {title}", font_size=16, color=ACCENT_CYAN, bold=True)
    add_text_box(slide, x, Inches(5.7), Inches(2.8), Inches(0.8),
                 desc, font_size=14, color=LIGHT_GRAY)

# ===== 第8页：实战项目 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.2), ACCENT_CYAN)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(8), Inches(0.6),
             "实战项目：智能个人助手 Agent", font_size=32, color=WHITE, bold=True)

# 功能模块
features = [
    ("📅 日程管理", "创建/查询/修改\n日历事件"),
    ("📧 邮件处理", "收发邮件、摘要\n自动回复"),
    ("🌤️ 信息查询", "天气、新闻\n百科、股票"),
    ("📁 文件操作", "读取/写入\n数据整理"),
    ("💬 对话记忆", "记住用户偏好\n跨会话上下文"),
]

for i, (title, desc) in enumerate(features):
    x = Inches(0.4) + Inches(i * 2.55)
    add_shape(slide, x, Inches(1.6), Inches(2.3), Inches(2.0), ACCENT_CARD)
    add_text_box(slide, x + Inches(0.15), Inches(1.8), Inches(2.0), Inches(0.4),
                 title, font_size=16, color=ACCENT_CYAN, bold=True)
    add_text_box(slide, x + Inches(0.15), Inches(2.3), Inches(2.0), Inches(1.0),
                 desc, font_size=14, color=LIGHT_GRAY)

# 技术栈
add_shape(slide, Inches(0.4), Inches(3.9), Inches(12.5), Inches(3.0), ACCENT_CARD)
add_text_box(slide, Inches(0.8), Inches(4.1), Inches(5), Inches(0.4),
             "🛠️ 推荐技术栈", font_size=20, color=YELLOW, bold=True)

techs = [
    ("框架", "LangChain / AutoGen\nCrewAI / Semantic Kernel"),
    ("LLM", "GPT-4o / Claude\nQwen / DeepSeek"),
    ("记忆", "ChromaDB / Pinecone\nRedis / PostgreSQL"),
    ("工具", "Function Calling\nREST API / Code Interpreter"),
    ("部署", "FastAPI + Docker\nKubernetes / Serverless"),
]

for i, (title, content) in enumerate(techs):
    x = Inches(0.8) + Inches(i * 2.45)
    add_text_box(slide, x, Inches(4.6), Inches(2.2), Inches(0.4),
                 title, font_size=14, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, x, Inches(5.1), Inches(2.2), Inches(1.2),
                 content, font_size=13, color=LIGHT_GRAY)

# ===== 第9页：最佳实践与挑战 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_line(slide, Inches(0.8), Inches(0.6), Inches(1.2), ACCENT_CYAN)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(8), Inches(0.6),
             "最佳实践与挑战", font_size=32, color=WHITE, bold=True)

# 左侧：最佳实践
add_shape(slide, Inches(0.8), Inches(1.6), Inches(5.8), Inches(5.2), ACCENT_CARD)
add_text_box(slide, Inches(1.2), Inches(1.8), Inches(5), Inches(0.4),
             "✅ 最佳实践", font_size=22, color=ACCENT_GREEN, bold=True)
add_bullet_text(slide, Inches(1.2), Inches(2.4), Inches(5), Inches(4.0),
                ["明确 Agent 的目标和边界", "设计优雅的错误处理机制",
                 "使用结构化输出（JSON Schema）", "实现人机协同（Human-in-the-Loop）",
                 "完善的日志与可观测性", "渐进式工具注册与权限管理",
                 "缓存与限流保护 LLM 调用", "定期评估 Agent 性能表现"],
                font_size=15, color=LIGHT_GRAY, spacing=Pt(6))

# 右侧：挑战
add_shape(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(5.2), ACCENT_CARD)
add_text_box(slide, Inches(7.4), Inches(1.8), Inches(4.8), Inches(0.4),
             "⚠️ 主要挑战", font_size=22, color=YELLOW, bold=True)
add_bullet_text(slide, Inches(7.4), Inches(2.4), Inches(4.8), Inches(4.0),
                ["幻觉控制：工具调用参数错误", "安全边界：防止提示注入攻击",
                 "成本控制：多轮调用 Token 消耗", "延迟优化：端到端响应时间",
                 "调试困难：思维链不可见", "鲁棒性：边缘情况处理",
                 "评估难：缺乏统一评测标准", "可解释性：Agent 决策透明度"],
                font_size=15, color=LIGHT_GRAY, spacing=Pt(6))

# ===== 第10页：结尾页 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# 中央内容
add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.0),
             "🚀 开始构建你的第一个 Agent", font_size=44, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5.5), Inches(3.2), Inches(2.5), ACCENT_BLUE)

add_text_box(slide, Inches(1.5), Inches(3.6), Inches(10), Inches(0.8),
             "从简单开始，逐步迭代，让智能体真正为你工作",
             font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
             "Thank You · 感谢聆听",
             font_size=18, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

# 保存到桌面
desktop = os.path.expanduser("~/Desktop")
output_path = os.path.join(desktop, "AI_Agent_开发指南.pptx")
prs.save(output_path)
print(f"✅ PPT 已生成到桌面：{output_path}")
print(f"   共 {len(prs.slides)} 页幻灯片")
